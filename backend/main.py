import asyncio
import io
import json
import os
import re
import tempfile
import zipfile
from contextlib import asynccontextmanager
from pathlib import Path

from dotenv import load_dotenv

load_dotenv(Path(__file__).parent.parent / ".env")

import httpx
import yt_dlp
from fastapi import FastAPI, File, HTTPException, Query, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, Response, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from starlette.middleware.base import BaseHTTPMiddleware

STATIC_DIR = Path(__file__).parent / "static"


class AppEntry(BaseModel):
    id: str
    name: str
    description: str
    url: str
    icon: str
    tags: list[str]


class AppsResponse(BaseModel):
    apps: list[AppEntry]


@asynccontextmanager
async def lifespan(app: FastAPI):
    config_path = Path(__file__).parent / "apps.json"
    with open(config_path) as f:
        data = json.load(f)
    app.state.apps = [AppEntry(**entry) for entry in data["apps"]]
    yield


app = FastAPI(title="Tools by Matt", version="0.1.0", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class NoCacheStaticMiddleware(BaseHTTPMiddleware):
    """Short cache for non-hashed static files so Cloudflare doesn't serve stale content."""
    async def dispatch(self, request, call_next):
        response = await call_next(request)
        path = request.url.path
        # Vite-hashed assets (e.g. /assets/index-Dy88f-5y.js) can cache forever
        if path.startswith("/assets/"):
            response.headers["cache-control"] = "public, max-age=31536000, immutable"
        # Sub-app static assets — short cache, must revalidate
        elif path.startswith("/apps/"):
            response.headers["cache-control"] = "public, max-age=60, must-revalidate"
        return response


app.add_middleware(NoCacheStaticMiddleware)


# ---------- API ----------

@app.get("/api/apps", response_model=AppsResponse)
async def get_apps():
    return AppsResponse(apps=app.state.apps)


# Bidirectional synonym groups — any word in a group matches any other
_SYNONYM_GROUPS = [
    {"food", "meal", "restaurant", "dining", "eat", "lunch", "dinner"},
    {"photo", "picture", "image", "camera", "pic"},
    {"measure", "ruler", "scale", "size", "dimension", "length", "distance"},
    {"random", "spin", "wheel", "pick", "choose", "roulette"},
    {"location", "map", "place", "nearby", "local", "drive"},
    {"perspective", "3d", "angle", "depth", "height"},
    {"youtube", "video", "audio", "music", "mp3", "mp4", "download", "dj", "mixer", "song"},
    {"file", "convert", "powerpoint", "pptx", "pdf", "image", "extract", "slides", "presentation", "video", "audio"},
]

_SYNONYMS: dict[str, set[str]] = {}
for group in _SYNONYM_GROUPS:
    for word in group:
        _SYNONYMS.setdefault(word, set()).update(group)


def _expand_query(query: str) -> set[str]:
    """Return the query plus any synonyms for words in it."""
    terms = {query}
    for word in query.split():
        terms.update(_SYNONYMS.get(word, set()))
    return terms


@app.get("/api/apps/search", response_model=AppsResponse)
async def search_apps(q: str = Query("", min_length=0)):
    if not q:
        return AppsResponse(apps=app.state.apps)
    terms = _expand_query(q.lower())
    results = [
        a for a in app.state.apps
        if any(
            t in a.name.lower()
            or t in a.description.lower()
            or any(t in tag for tag in a.tags)
            for t in terms
        )
    ]
    return AppsResponse(apps=results)


# ---------- Isochrone / Geocode Proxy (OpenRouteService) ----------

ORS_KEY = os.environ.get("ORS_API_KEY", "")


@app.get("/api/isochrone")
async def isochrone(lat: float, lon: float, minutes: int = Query(ge=1, le=60)):
    if not ORS_KEY:
        raise HTTPException(503, "ORS_API_KEY not configured")
    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.post(
            "https://api.openrouteservice.org/v2/isochrones/driving-car",
            headers={"Authorization": ORS_KEY},
            json={"locations": [[lon, lat]], "range": [minutes * 60], "range_type": "time"},
        )
    if resp.status_code != 200:
        raise HTTPException(resp.status_code, resp.text)
    return resp.json()


@app.get("/api/geocode")
async def geocode(q: str = Query(min_length=1)):
    if not ORS_KEY:
        raise HTTPException(503, "ORS_API_KEY not configured")
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.get(
            "https://api.openrouteservice.org/geocode/search",
            headers={"Authorization": ORS_KEY},
            params={"text": q, "size": 5},
        )
    if resp.status_code != 200:
        raise HTTPException(resp.status_code, resp.text)
    return resp.json()


# ---------- YouTube Downloader ----------

_YT_URL_RE = re.compile(r"(youtube\.com|youtu\.be)")

WAVEFORM_BARS = 200


def _generate_waveform(video_path: str, bars: int = WAVEFORM_BARS) -> list[float]:
    """Extract audio from video with ffmpeg and compute peak waveform."""
    import struct
    tmp_wav = video_path + ".wav"
    try:
        subprocess.run(
            [shutil.which("ffmpeg") or "ffmpeg", "-i", video_path,
             "-vn", "-ac", "1", "-ar", "8000", "-f", "s16le", "-y", tmp_wav],
            capture_output=True, timeout=30,
        )
        if not os.path.exists(tmp_wav):
            return []
        with open(tmp_wav, "rb") as f:
            raw = f.read()
        samples = struct.unpack(f"<{len(raw)//2}h", raw)
        samples_per_bar = max(1, len(samples) // bars)
        peaks: list[float] = []
        for i in range(bars):
            start = i * samples_per_bar
            end = min(start + samples_per_bar, len(samples))
            peak = max(abs(s) for s in samples[start:end]) if start < end else 0
            peaks.append(round(peak / 32768, 4))
        return peaks
    finally:
        if os.path.exists(tmp_wav):
            os.remove(tmp_wav)


@app.get("/api/yt/info")
async def yt_info(url: str = Query(...)):
    if not _YT_URL_RE.search(url):
        raise HTTPException(400, "Only YouTube URLs are supported")

    opts = {"quiet": True, "no_warnings": True, "skip_download": True}

    def extract():
        with yt_dlp.YoutubeDL(opts) as ydl:
            return ydl.extract_info(url, download=False)

    try:
        info = await asyncio.to_thread(extract)
    except yt_dlp.utils.DownloadError as e:
        raise HTTPException(400, str(e))

    return {
        "title": info.get("title"),
        "thumbnail": info.get("thumbnail"),
        "duration": info.get("duration"),
        "id": info.get("id"),
    }


@app.get("/api/yt/download")
async def yt_download(
    url: str = Query(...),
    format: str = Query("mp3", pattern="^(mp3|mp4)$"),
    quality: int = Query(1080, ge=360, le=2160),
):
    if not _YT_URL_RE.search(url):
        raise HTTPException(400, "Only YouTube URLs are supported")

    tmp_dir = tempfile.mkdtemp()
    out_template = os.path.join(tmp_dir, "%(title)s.%(ext)s")

    base_opts = {
        "quiet": True,
        "no_warnings": True,
        "outtmpl": out_template,
        "retries": 3,
        "fragment_retries": 3,
        "socket_timeout": 30,
    }

    if format == "mp3":
        opts = {
            **base_opts,
            "format": "bestaudio/best",
            "postprocessors": [
                {"key": "FFmpegExtractAudio", "preferredcodec": "mp3", "preferredquality": "192"}
            ],
        }
    else:
        opts = {
            **base_opts,
            "format": f"bestvideo[height<={quality}][vcodec^=avc1]+bestaudio[acodec^=mp4a]/bestvideo[height<={quality}]+bestaudio/best[height<={quality}]",
            "merge_output_format": "mp4",
            "postprocessors": [
                {"key": "FFmpegVideoConvertor", "preferedformat": "mp4"},
            ],
        }

    def download():
        with yt_dlp.YoutubeDL(opts) as ydl:
            ydl.download([url])

    try:
        await asyncio.to_thread(download)
    except yt_dlp.utils.DownloadError as e:
        raise HTTPException(400, str(e))

    # Find the output file
    files = list(Path(tmp_dir).iterdir())
    if not files:
        raise HTTPException(500, "Download produced no output")
    out_file = files[0]

    def stream():
        try:
            with open(out_file, "rb") as f:
                while chunk := f.read(65536):
                    yield chunk
        finally:
            import shutil
            shutil.rmtree(tmp_dir, ignore_errors=True)

    media_type = "audio/mpeg" if format == "mp3" else "video/mp4"
    filename = out_file.name
    return StreamingResponse(
        stream(),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.get("/api/yt/waveform")
async def yt_waveform(url: str = Query(...), bars: int = Query(200, ge=50, le=500)):
    if not _YT_URL_RE.search(url):
        raise HTTPException(400, "Only YouTube URLs are supported")

    tmp_dir = tempfile.mkdtemp()
    out_template = os.path.join(tmp_dir, "%(title)s.%(ext)s")
    opts = {
        "quiet": True,
        "no_warnings": True,
        "format": "bestaudio[filesize<20M]/bestaudio/best",
        "outtmpl": out_template,
        "retries": 3,
        "fragment_retries": 3,
        "socket_timeout": 30,
        "postprocessors": [
            {"key": "FFmpegExtractAudio", "preferredcodec": "wav", "preferredquality": "0"}
        ],
    }

    def run():
        with yt_dlp.YoutubeDL(opts) as ydl:
            ydl.download([url])
        files = list(Path(tmp_dir).iterdir())
        if not files:
            return []
        return _generate_waveform(str(files[0]), bars)

    try:
        peaks = await asyncio.to_thread(run)
    except Exception as e:
        raise HTTPException(400, str(e))
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

    return {"peaks": peaks}


# ---------- File Converter ----------

import shutil
import subprocess

IMAGE_EXTS = {"png", "jpg", "jpeg", "webp", "gif", "bmp", "tiff", "heic", "svg"}
VIDEO_EXTS = {"mp4", "mov", "avi", "mkv", "webm", "flv", "wmv", "m4v", "ts", "mts", "3gp"}
AUDIO_EXTS = {"mp3", "wav", "aac", "flac", "ogg", "m4a", "wma", "opus"}

# Extensions LibreOffice can convert to PDF
_LIBREOFFICE_EXTS = {
    "doc", "docx", "odt", "rtf", "txt",                     # word processors
    "xls", "xlsx", "ods", "csv",                             # spreadsheets
    "ppt", "pptx", "odp",                                    # presentations
    "pages", "numbers", "key",                               # iWork (macOS LO)
    "html", "htm",                                           # web
}

_SOFFICE = shutil.which("libreoffice") or shutil.which("soffice")
_FFMPEG = shutil.which("ffmpeg")


def _get_output_options(ext: str) -> list[str]:
    """Return available output formats for a given input extension."""
    if ext in ("pptx", "ppt", "odp", "key"):
        return ["images", "pdf"]
    if ext in ("doc", "docx", "odt", "rtf", "txt", "pages", "html", "htm"):
        return ["pdf", "images"]
    if ext in ("xls", "xlsx", "ods", "csv", "numbers"):
        return ["pdf", "images"]
    if ext == "pdf":
        return ["images"]
    if ext in IMAGE_EXTS:
        return ["png", "jpg", "webp", "pdf"]
    if ext in VIDEO_EXTS:
        opts = ["mp4", "webm", "mov", "mp3", "wav", "gif"]
        # Don't offer the same format as input
        return [o for o in opts if o != ext]
    if ext in AUDIO_EXTS:
        opts = ["mp3", "wav", "aac", "flac", "ogg"]
        return [o for o in opts if o != ext]
    # Unknown — try LibreOffice if available
    if _SOFFICE:
        return ["pdf", "images"]
    return []


def _extract_images_pptx(data: bytes) -> list[tuple[str, bytes]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    images: list[tuple[str, bytes]] = []
    idx = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                ext = img.content_type.split("/")[-1].replace("jpeg", "jpg")
                idx += 1
                images.append((f"slide{slide_num:02d}_{idx:03d}.{ext}", img.blob))
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    if hasattr(s, "image"):
                        img = s.image
                        ext = img.content_type.split("/")[-1].replace("jpeg", "jpg")
                        idx += 1
                        images.append((f"slide{slide_num:02d}_{idx:03d}.{ext}", img.blob))
    return images


def _extract_images_docx(data: bytes) -> list[tuple[str, bytes]]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    images: list[tuple[str, bytes]] = []
    for i, rel in enumerate(doc.part.rels.values(), 1):
        if "image" in rel.reltype:
            blob = rel.target_part.blob
            ct = rel.target_part.content_type
            ext = ct.split("/")[-1].replace("jpeg", "jpg")
            images.append((f"image_{i:03d}.{ext}", blob))
    return images


def _render_pdf_pages(data: bytes) -> list[tuple[str, bytes]]:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    images: list[tuple[str, bytes]] = []
    for page_num, page in enumerate(doc, 1):
        pix = page.get_pixmap(dpi=200)
        images.append((f"page{page_num:03d}.png", pix.tobytes("png")))
    return images


def _convert_to_pdf_libreoffice(data: bytes, filename: str) -> bytes:
    if not _SOFFICE:
        raise RuntimeError("LibreOffice not installed")
    tmp = tempfile.mkdtemp()
    try:
        src = os.path.join(tmp, filename)
        with open(src, "wb") as f:
            f.write(data)
        subprocess.run(
            [_SOFFICE, "--headless", "--convert-to", "pdf", "--outdir", tmp, src],
            capture_output=True, timeout=60,
        )
        stem = filename.rsplit(".", 1)[0]
        pdf_path = os.path.join(tmp, f"{stem}.pdf")
        if not os.path.exists(pdf_path):
            raise RuntimeError("LibreOffice conversion produced no output")
        with open(pdf_path, "rb") as f:
            return f.read()
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def _convert_with_ffmpeg(data: bytes, in_ext: str, out_ext: str, filename: str) -> tuple[str, bytes]:
    """Convert video/audio using ffmpeg."""
    if not _FFMPEG:
        raise RuntimeError("ffmpeg not installed")
    tmp = tempfile.mkdtemp()
    try:
        stem = filename.rsplit(".", 1)[0] if "." in filename else filename
        src = os.path.join(tmp, f"input.{in_ext}")
        out = os.path.join(tmp, f"{stem}.{out_ext}")
        with open(src, "wb") as f:
            f.write(data)

        cmd = [_FFMPEG, "-i", src, "-y"]
        if out_ext == "gif":
            # Video → GIF: scale down, 10fps, good quality palette
            cmd += ["-vf", "fps=10,scale=480:-1:flags=lanczos", "-loop", "0"]
        elif out_ext in ("mp3", "wav", "aac", "flac", "ogg"):
            # Extract/convert audio only
            if out_ext == "mp3":
                cmd += ["-vn", "-acodec", "libmp3lame", "-q:a", "2"]
            elif out_ext == "wav":
                cmd += ["-vn", "-acodec", "pcm_s16le"]
            elif out_ext == "aac":
                cmd += ["-vn", "-acodec", "aac", "-b:a", "192k"]
            elif out_ext == "flac":
                cmd += ["-vn", "-acodec", "flac"]
            elif out_ext == "ogg":
                cmd += ["-vn", "-acodec", "libvorbis", "-q:a", "5"]
        elif out_ext == "mp4":
            cmd += ["-c:v", "libx264", "-preset", "fast", "-crf", "23",
                    "-c:a", "aac", "-b:a", "192k", "-movflags", "+faststart"]
        elif out_ext == "webm":
            cmd += ["-c:v", "libvpx-vp9", "-crf", "30", "-b:v", "0",
                    "-c:a", "libopus", "-b:a", "128k"]
        elif out_ext == "mov":
            cmd += ["-c:v", "libx264", "-c:a", "aac", "-b:a", "192k"]
        cmd.append(out)

        result = subprocess.run(cmd, capture_output=True, timeout=300)
        if not os.path.exists(out):
            err = result.stderr.decode(errors="replace")[-500:]
            raise RuntimeError(f"ffmpeg failed: {err}")
        with open(out, "rb") as f:
            return f"{stem}.{out_ext}", f.read()
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def _convert_image(data: bytes, target_format: str) -> tuple[str, bytes]:
    from PIL import Image

    try:
        img = Image.open(io.BytesIO(data))
        if target_format.lower() == "pdf":
            if img.mode == "RGBA":
                img = img.convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="PDF")
            return "converted.pdf", buf.getvalue()
        if img.mode == "RGBA" and target_format.lower() in ("jpg", "jpeg"):
            img = img.convert("RGB")
        buf = io.BytesIO()
        fmt_map = {"jpg": "JPEG", "jpeg": "JPEG", "png": "PNG", "webp": "WEBP", "gif": "GIF"}
        img.save(buf, format=fmt_map.get(target_format.lower(), "PNG"))
        return f"converted.{target_format}", buf.getvalue()
    except Exception:
        # Pillow can't open it (HEIC, AVIF, etc.) — fall back to ffmpeg
        if not _FFMPEG:
            raise
        return _convert_with_ffmpeg(data, "img", target_format, f"image.{target_format}")


@app.get("/api/convert/formats")
async def convert_formats(ext: str = Query(...)):
    options = _get_output_options(ext.lower().lstrip("."))
    return {"options": options}


@app.post("/api/convert")
async def convert_file(
    file: UploadFile = File(...),
    target: str = Query("images"),
):
    data = await file.read()
    filename = file.filename or "file"
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    results: list[tuple[str, bytes]] = []

    # --- Video / Audio (ffmpeg) ---
    if ext in VIDEO_EXTS or ext in AUDIO_EXTS:
        name, converted = await asyncio.to_thread(
            _convert_with_ffmpeg, data, ext, target, filename
        )
        results = [(name, converted)]

    # --- Images ---
    elif ext in IMAGE_EXTS:
        if target in ("png", "jpg", "webp", "pdf"):
            name, converted = await asyncio.to_thread(_convert_image, data, target)
            results = [(name, converted)]
        else:
            results = [(filename, data)]

    # --- PowerPoint (direct image extraction) ---
    elif ext in ("pptx",) and target == "images":
        results = await asyncio.to_thread(_extract_images_pptx, data)

    # --- Word (direct image extraction) ---
    elif ext in ("docx",) and target == "images":
        results = await asyncio.to_thread(_extract_images_docx, data)

    # --- PDF ---
    elif ext == "pdf":
        results = await asyncio.to_thread(_render_pdf_pages, data)

    # --- Document → PDF (via LibreOffice) ---
    elif target == "pdf" and ext in (_LIBREOFFICE_EXTS | {"pptx", "ppt", "docx", "doc"}):
        pdf_data = await asyncio.to_thread(_convert_to_pdf_libreoffice, data, filename)
        results = [(filename.rsplit(".", 1)[0] + ".pdf", pdf_data)]

    # --- Document → images (LibreOffice → PDF → render pages) ---
    elif target == "images" and (ext in _LIBREOFFICE_EXTS or ext in ("ppt", "doc", "odp", "key", "pages")):
        pdf_data = await asyncio.to_thread(_convert_to_pdf_libreoffice, data, filename)
        results = await asyncio.to_thread(_render_pdf_pages, pdf_data)

    # --- Unknown: try LibreOffice as fallback ---
    elif _SOFFICE:
        try:
            pdf_data = await asyncio.to_thread(_convert_to_pdf_libreoffice, data, filename)
            if target == "images":
                results = await asyncio.to_thread(_render_pdf_pages, pdf_data)
            else:
                results = [(filename.rsplit(".", 1)[0] + ".pdf", pdf_data)]
        except Exception:
            raise HTTPException(400, f"Cannot convert .{ext} files")
    else:
        raise HTTPException(400, f"Cannot convert .{ext} files (install LibreOffice for more formats)")

    if not results:
        raise HTTPException(400, "No content could be extracted")

    # Single file: return directly
    if len(results) == 1:
        name, content = results[0]
        ext_out = name.rsplit(".", 1)[-1]
        media_types = {
            "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "gif": "image/gif", "webp": "image/webp", "pdf": "application/pdf",
            "mp4": "video/mp4", "webm": "video/webm", "mov": "video/quicktime",
            "mp3": "audio/mpeg", "wav": "audio/wav", "aac": "audio/aac",
            "flac": "audio/flac", "ogg": "audio/ogg",
        }
        return Response(
            content,
            media_type=media_types.get(ext_out, "application/octet-stream"),
            headers={"Content-Disposition": f'attachment; filename="{name}"'},
        )

    # Multiple files: zip them
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, content in results:
            zf.writestr(name, content)
    buf.seek(0)
    stem = filename.rsplit(".", 1)[0]
    return Response(
        buf.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{stem}_converted.zip"'},
    )


# ---------- Static files (production: built frontend served by FastAPI) ----------

# React build assets (/assets/*)
assets_dir = STATIC_DIR / "assets"
if assets_dir.is_dir():
    app.mount("/assets", StaticFiles(directory=str(assets_dir)), name="assets")

# SPA fallback — serve index.html for any non-API, non-static path
index_html = STATIC_DIR / "index.html"


@app.get("/{path:path}")
async def spa_fallback(request: Request, path: str):
    # Try to serve the exact file first (e.g. favicon.ico, app JS/CSS assets)
    file_path = STATIC_DIR / path
    if path and file_path.is_file():
        return FileResponse(file_path)
    # Otherwise serve the React SPA (client-side routing handles /apps/*)
    if index_html.is_file():
        return FileResponse(index_html)
    return {"detail": "Not found — run `make build` to generate frontend"}
